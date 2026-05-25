"""
Generate the CSCI435 defence slide deck.

Output: ``slides/SignSpeak_AI_Defence.pptx``

Twelve slides matching the report structure, with embedded charts from
``ml/results/`` and a dark professional theme.
"""

from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "SignSpeak_AI_Defence.pptx"
FIG_DIR = ROOT / "report" / "figures"
RESULTS_DIR = ROOT / "ml" / "results"


BG = RGBColor(0x0B, 0x12, 0x20)
FG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT = RGBColor(0x60, 0xA5, 0xFA)
SUBTLE = RGBColor(0x94, 0xA3, 0xB8)


def paint_background(slide, bg=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg


def add_text(slide, x_cm, y_cm, w_cm, h_cm, text, *, size=18, bold=False,
             color=FG, align="left"):
    tx = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_accent_bar(slide, y_cm: float = 1.6):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2.5), Cm(y_cm), Cm(2.2), Cm(0.18))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT


def add_image_if_exists(slide, path: Path, x_cm, y_cm, w_cm):
    if not path.exists():
        return False
    slide.shapes.add_picture(str(path), Cm(x_cm), Cm(y_cm), width=Cm(w_cm))
    return True


def read_csv(path: Path) -> list[list[str]]:
    if not path.exists():
        return []
    with path.open(encoding="utf-8", newline="") as f:
        return list(csv.reader(f))


def add_table(slide, x_cm, y_cm, w_cm, h_cm, rows: list[list[str]]):
    if not rows:
        add_text(slide, x_cm, y_cm, w_cm, h_cm,
                 "[table to be inserted after training]",
                 size=14, color=SUBTLE)
        return
    table = slide.shapes.add_table(
        rows=len(rows), cols=len(rows[0]),
        left=Cm(x_cm), top=Cm(y_cm),
        width=Cm(w_cm), height=Cm(h_cm),
    ).table
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B) if i == 0 else RGBColor(0x0F, 0x17, 0x2A)
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(cell)
            r.font.size = Pt(11)
            r.font.color.rgb = FG
            r.font.bold = (i == 0)


# ---------- slides ----------


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    paint_background(s)
    add_text(s, 2.5, 5.0, 24, 2.0, "SignSpeak AI", size=60, bold=True, align="center")
    add_text(s, 2.5, 8.4, 24, 1.2,
             "Real-Time American + Arabic Sign Language Translation",
             size=22, color=ACCENT, align="center")
    add_text(s, 2.5, 11.0, 24, 1.0, "CSCI435 - University of Wollongong in Dubai",
             size=18, color=SUBTLE, align="center")
    add_text(s, 2.5, 12.4, 24, 1.0,
             "Yousef Kanjo, Member 2, Member 3, Member 4, Member 5",
             size=14, color=SUBTLE, align="center")


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Problem and user story", size=32, bold=True)
    add_text(
        s, 2.5, 3.2, 24, 8.0,
        "- Deaf UOWD students rely on a small pool of sign-language "
        "interpreters whose availability is constrained.\n\n"
        "- When an interpreter is unavailable, normal office-hour and "
        "lecture conversations break down.\n\n"
        "- Aisha (deaf, 3rd-year CS) opens SignSpeak AI in any browser. "
        "She signs in Arabic Sign Language; the system shows the text "
        "on screen and reads it aloud through Web Speech TTS. The hearing "
        "professor responds verbally. Conversation flows.",
        size=18,
    )


def slide_capabilities(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Eight CV capabilities (only 4 required)",
             size=30, bold=True)
    items = [
        "1. Image enhancement - CLAHE",
        "2. Image segmentation - MediaPipe Selfie",
        "3. Binary morphological ops - erode + dilate",
        "4. Keypoint detection - MediaPipe Holistic",
        "5. Face detection - 468 face landmarks",
        "6. Edge detection - Canny on hand crop",
        "7. Video processing - optical flow + voting",
        "8. Object recognition - MobileNetV3 + MLP ensemble",
    ]
    add_text(s, 2.5, 3.0, 24, 9.5, "\n\n".join(items), size=18)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "System architecture", size=30, bold=True)
    if not add_image_if_exists(s, FIG_DIR / "architecture.png", 2.0, 3.0, 25.0):
        add_text(s, 2.5, 3.0, 24, 8.0,
                 "[Architecture diagram - regenerate via report/generate_architecture_diagram.py]",
                 size=16, color=SUBTLE)


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Per-frame pipeline", size=30, bold=True)
    add_text(
        s, 2.5, 3.0, 24, 10,
        "Each WebSocket frame at 15 FPS:\n\n"
        "1. CLAHE enhancement on LAB L channel\n"
        "2. MediaPipe Selfie segmentation -> foreground mask\n"
        "3. Open + close morphology cleans the mask\n"
        "4. MediaPipe Holistic -> 21 hand + 468 face + 33 pose keypoints\n"
        "5. Canny edges on the hand crop (auxiliary)\n"
        "6. Optical-flow magnitude on consecutive hand landmarks\n"
        "7. Landmark MLP inference (63-D vector -> 61 classes)\n"
        "8. MobileNetV3-Small inference (96x96 crop -> 61 classes)\n"
        "    Ensemble + 8-of-10 vote + conf>=0.85 -> stable label",
        size=16,
    )


def slide_models(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Models and training", size=30, bold=True)
    add_text(
        s, 2.5, 3.0, 24, 9,
        "Landmark MLP (trained from scratch on our curated data)\n"
        "  - 63 -> 256 -> 128 -> 64 -> 61, BatchNorm, dropout 0.3\n"
        "  - 20 epochs, AdamW 1e-3, cosine annealing, label smoothing 0.05\n"
        "  - On-line augmentation: rotation +/-10 deg, scale jitter, noise\n\n"
        "MobileNetV3-Small (ImageNet pretrained, fine-tuned)\n"
        "  - Backbone frozen except last conv block\n"
        "  - Final classifier replaced with 1024 -> 61 linear\n"
        "  - 5 epochs, AdamW, ColorJitter + small rotations\n\n"
        "Both checkpoints loaded at startup by backend/app/models/ensemble.py.",
        size=16,
    )


def slide_results_accuracy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Results - accuracy", size=30, bold=True)
    rows = read_csv(RESULTS_DIR / "accuracy_table.csv")
    add_table(s, 2.5, 3.5, 22, 4.0, rows)
    add_image_if_exists(s, RESULTS_DIR / "training_curves.png", 2.5, 8.0, 22)


def slide_results_confusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Results - confusion matrices", size=30, bold=True)
    add_image_if_exists(s, RESULTS_DIR / "confusion_matrix_asl.png", 1.0, 2.8, 14)
    add_image_if_exists(s, RESULTS_DIR / "confusion_matrix_arsl.png", 15.5, 2.8, 14)


def slide_results_perf(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Results - latency, ablation", size=30, bold=True)
    add_text(s, 2.5, 3.0, 12, 1.0, "Latency & FPS", size=18, bold=True)
    add_table(s, 2.5, 4.0, 11, 4.0, read_csv(RESULTS_DIR / "fps_latency.csv"))
    add_text(s, 15.5, 3.0, 12, 1.0, "Ablation study", size=18, bold=True)
    add_table(s, 15.5, 4.0, 11, 4.0, read_csv(RESULTS_DIR / "ablation_table.csv"))


def slide_demo_plan(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Live demo plan", size=30, bold=True)
    add_text(
        s, 2.5, 3.0, 24, 10,
        "1. Open http://localhost:3000 and click 'Start live mode'\n"
        "2. Show webcam + skeleton overlay + 25 FPS in stats panel\n"
        "3. Sign A, B, C in ASL - letters appear; confidence > 95%\n"
        "4. Spell H-E-L-L-O - sentence builder fills\n"
        "5. Click 'Speak' - English TTS says 'Hello'\n"
        "6. Toggle to ArSL - sign 2 Arabic letters - Arabic glyphs appear\n"
        "7. Switch to 'Upload' page - drop a pre-recorded clip - transcript renders\n"
        "8. Back to live - dim the lights to show CLAHE robustness\n"
        "9. Q&A",
        size=16,
    )


def slide_limitations(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_accent_bar(s)
    add_text(s, 2.5, 1.0, 24, 1.4, "Limitations and future work", size=30, bold=True)
    add_text(
        s, 2.5, 3.0, 24, 10,
        "Limitations\n"
        "  - Fingerspelling only; no continuous sentence-level signs yet\n"
        "  - Single-hand pipeline; two-handed signs are not modelled\n"
        "  - 32 ArSL letters supported; ArSL has > 100 dynamic signs\n\n"
        "Future work\n"
        "  - Add WLASL temporal model for whole-word signs\n"
        "  - React Native client with TFLite for on-device inference\n"
        "  - Bidirectional - text-to-sign avatar for hearing -> deaf direction\n"
        "  - Expand custom dataset with UAE community contributions",
        size=16,
    )


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(s)
    add_text(s, 2.5, 5.0, 24, 2.0, "Thank you", size=60, bold=True, align="center")
    add_text(s, 2.5, 9.0, 24, 1.5,
             "Questions?", size=28, color=ACCENT, align="center")
    add_text(s, 2.5, 11.5, 24, 1.0,
             "github.com/ykanjo04/SignSpeak-AI",
             size=18, color=SUBTLE, align="center")


# ---------- main ----------


def main() -> None:
    prs = Presentation()
    prs.slide_width = Cm(33.867)         # 16:9
    prs.slide_height = Cm(19.05)

    slide_title(prs)
    slide_problem(prs)
    slide_capabilities(prs)
    slide_architecture(prs)
    slide_pipeline(prs)
    slide_models(prs)
    slide_results_accuracy(prs)
    slide_results_confusion(prs)
    slide_results_perf(prs)
    slide_demo_plan(prs)
    slide_limitations(prs)
    slide_thanks(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    main()
